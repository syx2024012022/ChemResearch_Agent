from __future__ import annotations

import hashlib
import json
import math
import textwrap
from pathlib import Path
from uuid import UUID

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from chemresearch_agent.domain.models import PresentationArtifact, SlideContent
from chemresearch_agent.skills.templates import BuiltinTemplateRegistry

WIDTH, HEIGHT, DPI = 1280, 720, 96


class PythonPptxPresentationRenderer:
    """Public-dependency fallback for environments without Artifact Tool."""

    version = "python-pptx-1.0"

    def __init__(self) -> None:
        self._templates = BuiltinTemplateRegistry()
        root = Path(__file__).resolve().parents[4] / "assets" / "templates" / "chem_group_standard"
        self._backgrounds = [root / "slide-1.png", root / "slide-2.png"]

    def render(self, session_id: UUID, contents: list[SlideContent], output_root: Path):
        self._warnings: list[str] = []
        payload = json.dumps(
            [item.model_dump(mode="json") for item in contents], sort_keys=True, ensure_ascii=False
        )
        digest = hashlib.sha256(f"{session_id}:{self.version}:{payload}".encode()).hexdigest()
        artifact_id = UUID(digest[:32])
        output = output_root / str(artifact_id)
        output.mkdir(parents=True, exist_ok=True)
        deck = Presentation()
        deck.slide_width, deck.slide_height = _emu(WIDTH), _emu(HEIGHT)
        blank = deck.slide_layouts[6]
        previews, layouts = [], []
        for number, content in enumerate(contents, start=1):
            slide = deck.slides.add_slide(blank)
            canvas = self._background(number)
            elements = self._compose(slide, canvas, content, number)
            slide.notes_slide.notes_text_frame.text = content.speaker_notes or "[Sources]\n- none"
            preview = output / f"slide-{number:02d}.png"
            canvas.save(preview)
            previews.append(str(preview.resolve()))
            layout = output / f"slide-{number:02d}.layout.json"
            layout.write_text(
                json.dumps({"elements": elements}, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            layouts.append(str(layout.resolve()))
        pptx = output / "presentation.pptx"
        deck.save(pptx)
        montage = output / "montage.webp"
        _montage([Path(item) for item in previews]).save(montage, "WEBP")
        return PresentationArtifact(
            artifact_id=artifact_id,
            session_id=session_id,
            pptx_path=str(pptx.resolve()),
            preview_paths=previews,
            layout_paths=layouts,
            montage_path=str(montage.resolve()),
            render_log=[
                "renderer=python-pptx",
                f"previews={len(previews)}",
                *dict.fromkeys(self._warnings),
            ],
            renderer_version=self.version,
            input_hash=digest,
            slide_count=len(contents),
        )

    def _background(self, number: int) -> Image.Image:
        path = self._backgrounds[0 if number == 1 else 1]
        return (
            Image.open(path).convert("RGB").resize((WIDTH, HEIGHT))
            if path.exists()
            else Image.new("RGB", (WIDTH, HEIGHT), "white")
        )

    def _compose(self, slide, canvas, content: SlideContent, number: int) -> list[dict]:
        background = self._backgrounds[0 if number == 1 else 1]
        if background.exists():
            slide.shapes.add_picture(str(background), 0, 0, width=_emu(WIDTH), height=_emu(HEIGHT))
        elements: list[dict] = []
        if number == 1:
            _rect(slide, canvas, (0, 144, 1280, 576), "FFFFFF")
            self._text(
                slide,
                canvas,
                elements,
                "title-1",
                content.title,
                (64, 154, 1152, 96),
                26,
                True,
                "102A2E",
                "center",
            )
            self._text(
                slide,
                canvas,
                elements,
                "title-authors",
                _text(content, "authors"),
                (70, 256, 1140, 36),
                22,
                False,
                "102A2E",
            )
            self._text(
                slide,
                canvas,
                elements,
                "title-publication",
                _text(content, "publication"),
                (70, 294, 1140, 32),
                17,
                False,
                "38536A",
            )
            toc = _assets(content, "toc_graphic")
            if toc:
                self._image(
                    slide, canvas, elements, "toc-graphic", toc[0].asset_path, (80, 334, 1120, 325)
                )
            return elements
        self._text(
            slide,
            canvas,
            elements,
            f"title-{number}",
            content.title,
            (64, 24, 1152, 72),
            36,
            True,
            "FFFFFF",
            "center",
        )
        message = _text(content, "message")
        images = [item for item in content.blocks if item.asset_path]
        layout = self._templates.get(content.template_id).layout
        for index, (block, box) in enumerate(
            zip(images, _image_boxes(layout, len(images)), strict=False), start=1
        ):
            self._image(slide, canvas, elements, f"image-{index}", block.asset_path, box)
        if message and message != content.title:
            box, style, circle = _message_box(layout, len(images))
            if circle:
                _circle(slide, canvas, circle)
                elements.append(_element(f"message-circle-{number}", circle))
            self._text(slide, canvas, elements, f"message-{number}", message, box, *style)
        _rect(slide, canvas, (1150, 668, 86, 38), "FFFFFF")
        self._text(
            slide,
            canvas,
            elements,
            f"page-{number}",
            str(number),
            (1160, 675, 56, 24),
            12,
            False,
            "6B7B7D",
            "right",
        )
        return elements

    def _text(self, slide, canvas, elements, name, text, box, size, bold, color, align="left"):
        if not text:
            return
        shape = slide.shapes.add_textbox(*map(_emu, box))
        shape.name = name
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
            align, PP_ALIGN.LEFT
        )
        font = paragraph.runs[0].font
        font.name = "SimHei" if _chinese(text) else "Arial"
        font.size, font.bold = Pt(size), bold
        font.color.rgb = RGBColor.from_string(color)
        _draw_text(canvas, text, box, size, bold, color, align, self._warnings)
        elements.append(
            {
                **_element(name, box),
                "text": text,
                "textLayout": {"lineCount": _lines(text, box[2], size)},
            }
        )

    def _image(self, slide, canvas, elements, name, asset, box):
        path = Path(asset) if asset else None
        if not path or not path.exists():
            return
        fitted = _fit(path, box)
        slide.shapes.add_picture(str(path), *map(_emu, fitted))
        image = ImageOps.contain(Image.open(path).convert("RGB"), (fitted[2], fitted[3]))
        canvas.paste(image, fitted[:2])
        elements.append(_element(name, fitted))


def _image_boxes(layout: str, count: int):
    if layout == "panel_triptych" and count >= 3:
        return [(48, 116, 550, 260), (48, 390, 550, 275), (620, 112, 610, 558)]
    if layout == "two_panels_fill" and count >= 2:
        return [(38, 112, 590, 565), (652, 112, 590, 565)]
    if layout == "weighted_two_images" and count >= 2:
        return [(38, 112, 525, 565), (582, 112, 660, 565)]
    if layout == "stacked_mechanism_overview" and count >= 2:
        return [(12, 82, 1256, 202), (250, 282, 780, 430)]
    if count >= 2:
        return [(64, 132, 550, 510), (666, 132, 550, 510)]
    if layout == "multipanel_full":
        return [(16, 102, 1248, 610)]
    if layout == "image_full":
        return [(80, 126, 1120, 530)]
    if layout == "image_left":
        return [(64, 132, 720, 520)]
    return [(445, 132, 775, 520)] if count else []


def _message_box(layout: str, count: int):
    if count >= 2:
        return (64, 600, 1152, 64), (18, False, "25383B", "left"), None
    if layout in {"image_full", "multipanel_full"}:
        return (80, 606, 1120, 62), (20, False, "25383B", "left"), None
    if layout == "image_left":
        return (790, 190, 410, 300), (20, False, "25383B", "left"), None
    if count == 1:
        return (102, 244, 292, 226), (26, True, "000000", "center"), (68, 180, 360, 360)
    return (120, 220, 1040, 240), (38, True, "1F55A5", "center"), None


def _rect(slide, canvas, box, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *map(_emu, box))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()
    ImageDraw.Draw(canvas).rectangle(
        (box[0], box[1], box[0] + box[2], box[1] + box[3]), fill=f"#{color}"
    )


def _circle(slide, canvas, box):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, *map(_emu, box))
    shape.fill.background()
    shape.line.color.rgb, shape.line.width = RGBColor(183, 183, 183), Pt(1.5)
    ImageDraw.Draw(canvas).ellipse(
        (box[0], box[1], box[0] + box[2], box[1] + box[3]), outline="#B7B7B7", width=2
    )


def _draw_text(canvas, text, box, size, bold, color, align, warnings):
    chinese = _chinese(text)
    font, requested_font_found = _font(size, bold, chinese)
    if chinese and not requested_font_found:
        warnings.append("cjk_preview_font_missing")
    draw = ImageDraw.Draw(canvas)
    x, y, width, height = box
    capacity = max(1, int(width / max(1, size * 0.52)))
    lines = textwrap.wrap(text, width=capacity, break_long_words=False) or [""]
    maximum_lines = max(1, int(height / max(1, size * 1.25)))
    if len(lines) > maximum_lines:
        lines = lines[:maximum_lines]
        lines[-1] = lines[-1].rstrip(" .") + "…"
    visible_text = "\n".join(lines)
    centered = align == "center"
    point, anchor = (
        ((x + width // 2, y + height // 2), "mm") if centered else ((x, y + height // 2), "lm")
    )
    draw.multiline_text(
        point,
        visible_text,
        font=font,
        fill=f"#{color}",
        anchor=anchor,
        align=align,
        spacing=4,
    )


def _font(size, bold, chinese=False):
    candidates = _font_candidates(chinese, bold)
    path = next((item for item in candidates if item.exists()), None)
    if path:
        return ImageFont.truetype(str(path), size), True
    fallback = _font_candidates(False, bold)
    path = next((item for item in fallback if item.exists()), None)
    font = ImageFont.truetype(str(path), size) if path else ImageFont.load_default()
    return font, not chinese


def _font_candidates(chinese, bold):
    if chinese:
        windows = [
            Path("C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc"),
            Path("C:/Windows/Fonts/simhei.ttf"),
        ]
        linux = [
            Path(
                "/usr/share/fonts/opentype/noto/"
                f"NotoSansCJK-{'Bold' if bold else 'Regular'}.ttc"
            ),
            Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
        ]
        return [*windows, *linux]
    return [
        Path(f"C:/Windows/Fonts/{'arialbd' if bold else 'arial'}.ttf"),
        Path(f"/usr/share/fonts/truetype/dejavu/DejaVuSans{'-Bold' if bold else ''}.ttf"),
    ]


def _fit(path: Path, box):
    x, y, width, height = box
    with Image.open(path) as image:
        scale = min(width / image.width, height / image.height)
        fitted = (round(image.width * scale), round(image.height * scale))
    return (x + (width - fitted[0]) // 2, y + (height - fitted[1]) // 2, *fitted)


def _text(content, slot):
    block = next((item for item in content.blocks if item.slot == slot), None)
    return block.text if block and block.text else ""


def _assets(content, slot):
    return [item for item in content.blocks if item.slot == slot and item.asset_path]


def _element(name, box):
    return {"name": name, "bbox": list(box)}


def _emu(pixels):
    return Inches(pixels / DPI)


def _chinese(text):
    return any("\u3400" <= item <= "\u9fff" for item in text)


def _lines(text, width, size):
    return max(1, math.ceil(len(text) / max(1, int(width / max(1, size * 0.52)))))


def _montage(paths):
    images = [Image.open(path).convert("RGB").resize((320, 180)) for path in paths]
    columns, rows = min(3, max(1, len(images))), math.ceil(len(images) / 3)
    result = Image.new("RGB", (columns * 320, rows * 180), "white")
    for index, image in enumerate(images):
        result.paste(image, ((index % columns) * 320, (index // columns) * 180))
    return result
