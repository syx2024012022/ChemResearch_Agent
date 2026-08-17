from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from PIL import Image
from pptx import Presentation

from chemresearch_agent.domain.models import ContentBlock, SlideContent
from chemresearch_agent.tools.presentation import PythonPptxPresentationRenderer
from chemresearch_agent.tools.presentation.python_pptx_renderer import _font_candidates


def test_public_fallback_renderer_creates_editable_deck_notes_and_previews(
    tmp_path: Path,
) -> None:
    figure = tmp_path / "figure.png"
    Image.new("RGB", (900, 420), "white").save(figure)
    contents = [
        SlideContent(
            slide_id="title",
            origin_plan_slide_id="title",
            title="Fallback Renderer Test",
            template_id="title_paper_toc",
            blocks=[
                ContentBlock(slot="message", text="Fallback Renderer Test"),
                ContentBlock(slot="authors", text="Test Author"),
                ContentBlock(slot="publication", text="Test Journal · 2026"),
                ContentBlock(slot="toc_graphic", asset_path=str(figure), figure_id="toc"),
            ],
            speaker_notes="[Sources]\n- Test metadata",
        ),
        SlideContent(
            slide_id="result",
            origin_plan_slide_id="result",
            title="Key Transformation",
            template_id="image_full",
            blocks=[
                ContentBlock(slot="message", text="A concise grounded conclusion"),
                ContentBlock(slot="figure_1", asset_path=str(figure), figure_id="figure-1"),
            ],
            speaker_notes="[Sources]\n- Figure 1, PDF p.2",
        ),
    ]
    artifact = PythonPptxPresentationRenderer().render(uuid4(), contents, tmp_path / "out")
    assert Path(artifact.pptx_path).is_file()
    assert len(artifact.preview_paths) == 2
    assert len(artifact.layout_paths) == 2
    assert Path(artifact.montage_path).is_file()
    deck = Presentation(artifact.pptx_path)
    assert len(deck.slides) == 2
    assert "Fallback Renderer Test" in deck.slides[0].shapes[2].text
    assert "[Sources]" in deck.slides[0].notes_slide.notes_text_frame.text
    assert "Figure 1" in deck.slides[1].notes_slide.notes_text_frame.text


def test_preview_font_candidates_include_windows_and_linux_cjk_fonts() -> None:
    names = {str(path).replace("\\", "/") for path in _font_candidates(True, False)}
    assert any(name.endswith("/Windows/Fonts/msyh.ttc") for name in names)
    assert any("NotoSansCJK-Regular.ttc" in name for name in names)
